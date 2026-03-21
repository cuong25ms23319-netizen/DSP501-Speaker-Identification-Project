"""
Generate Speaker Identification presentation as PPTX
using figures from Jupyter notebooks.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGURES = os.path.join(os.path.dirname(__file__), 'figures')
OUTPUT = os.path.join(os.path.dirname(__file__), 'DSP Study', 'Speaker_Identification_Slides.pptx')

# FPT Orange
FPT = RGBColor(0xF3, 0x70, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=FPT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=Pt(6), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def add_image(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width:
        kwargs['width'] = Inches(width)
    if height:
        kwargs['height'] = Inches(height)
    slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


def add_orange_bar(slide, top=0, height=1.2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top),
        prs.slide_width, Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = FPT
    shape.line.fill.background()


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, FPT)
add_textbox(slide, 0.5, 1.5, 12, 1.5, 'Speaker Identification System',
            font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 3.0, 12, 1, 'DSP-Enhanced Pipeline vs. Baseline',
            font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 4.5, 12, 0.8,
            'Nguyen Huy Cuong  |  Hon Vi Dan  |  Le Nhut Thanh Quang  |  Nguyen Duc Minh Khoa',
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 5.5, 12, 0.8,
            'FPT University — DSP501  |  Supervisor: Dr. Dang Ngoc Minh Duc',
            font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Outline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Outline', font_size=36, bold=True, color=WHITE)
items = [
    '1.  Introduction — Problem & Motivation',
    '2.  Dataset — 5 speakers, 125 samples',
    '3.  DSP Methodology — FIR Filter & Pre-emphasis',
    '4.  Feature Engineering — Pipeline A vs Pipeline B',
    '5.  AI Modeling — SVM with RBF Kernel',
    '6.  Experimental Results',
    '7.  Discussion — Why Pipeline B Wins',
    '8.  Limitations & Conclusion',
    '9.  Live Demo',
]
tf = add_textbox(slide, 1, 1.8, 11, 5, items[0], font_size=24, color=BLACK)
for item in items[1:]:
    add_para(tf, item, font_size=24, space_before=Pt(14))

# ============================================================
# SLIDE 3: Problem & Motivation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Problem & Motivation', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.6, 12, 1, 'Goal: Identify WHO is speaking from a short audio clip',
                 font_size=24, bold=True)
add_para(tf, '', font_size=12)
add_para(tf, 'Challenge:', font_size=22, bold=True, space_before=Pt(16))
add_para(tf, '  - Speech is non-stationary — properties change over time', font_size=20)
add_para(tf, '  - Environmental noise corrupts the signal', font_size=20)
add_para(tf, '  - Different mics/sessions produce different characteristics', font_size=20)
add_para(tf, '', font_size=12)
add_para(tf, 'Research Question:', font_size=22, bold=True, space_before=Pt(16))
add_para(tf, '  Does DSP preprocessing improve speaker identification accuracy?', font_size=20)
add_para(tf, '', font_size=12)
add_para(tf, 'Approach: Compare two pipelines', font_size=22, bold=True, space_before=Pt(16))
add_para(tf, '  Pipeline A:  Raw signal  →  basic time features (6-dim)  →  SVM', font_size=20, color=GRAY)
add_para(tf, '  Pipeline B:  FIR + Pre-emphasis  →  MFCC features (26-dim)  →  SVM', font_size=20, color=FPT, bold=True)

# ============================================================
# SLIDE 4: Dataset
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Dataset', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.6, 6, 4, '5 Speakers x 25 files = 125 samples', font_size=26, bold=True)
add_para(tf, '', font_size=10)
add_para(tf, 'Speaker 07 — Dan', font_size=20)
add_para(tf, 'Speaker 08 — Cuong', font_size=20)
add_para(tf, 'Speaker 09 — Quang', font_size=20)
add_para(tf, 'Speaker 10 — Anne', font_size=20)
add_para(tf, 'Speaker 11 — Khoa', font_size=20)
add_para(tf, '', font_size=10)
add_para(tf, 'Duration: ~3 seconds each', font_size=18, color=GRAY)
add_para(tf, 'Mono WAV, fs = 16,000 Hz', font_size=18, color=GRAY)
add_para(tf, '48,000 samples per clip', font_size=18, color=GRAY)

tf2 = add_textbox(slide, 7, 1.6, 5.5, 4, 'Preprocessing Pipeline:', font_size=24, bold=True)
add_para(tf2, '', font_size=10)
add_para(tf2, '1. Load audio as mono 16 kHz', font_size=20)
add_para(tf2, '2. Normalize amplitude to [-1, 1]', font_size=20)
add_para(tf2, '3. Trim silence (threshold: -20 dB)', font_size=20)
add_para(tf2, '4. Pad or crop to 48,000 samples', font_size=20)

# ============================================================
# SLIDE 5: Waveform — raw signal
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Signal Analysis — Waveform', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'waveform_comparison.png'), 0.5, 1.2, width=12)
tf = add_textbox(slide, 0.5, 5.5, 12, 1.5,
                 'Raw signal (top) vs Filtered signal (bottom) — FIR bandpass 300-3400 Hz removes noise while preserving speech',
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 6: FIR Filter Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'FIR Bandpass Filter (300–3400 Hz)', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'fir_analysis.png'), 0.3, 1.2, width=12.5)
tf = add_textbox(slide, 0.5, 5.3, 12, 2,
                 'Left: Magnitude response — passband at 0 dB, stopband < -40 dB\n'
                 'Center: Linear phase in passband — no signal distortion\n'
                 'Right: Symmetric impulse response (101 taps, Hamming window)',
                 font_size=16, color=GRAY)

# ============================================================
# SLIDE 7: Spectrum comparison (FFT)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Spectrum Comparison (FFT)', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'spectrum_comparison.png'), 0.5, 1.2, width=12)
tf = add_textbox(slide, 0.5, 5.8, 12, 1,
                 'Blue = Raw, Orange = Filtered — energy outside 300-3400 Hz is heavily attenuated',
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 8: STFT Spectrogram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Spectrogram (STFT) — Before vs After FIR', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'stft_comparison.png'), 0.3, 1.2, width=12.5)
tf = add_textbox(slide, 0.5, 5.8, 12, 1,
                 'Formants (bright horizontal bands) preserved in 300-3400 Hz; noise outside the band removed',
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 9: PSD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Power Spectral Density (PSD)', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'psd_comparison.png'), 0.8, 1.2, width=11.5)
tf = add_textbox(slide, 0.5, 5.8, 12, 1,
                 "Welch's method — power drops sharply outside 300-3400 Hz after filtering (red shaded = removed)",
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 10: SNR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Signal-to-Noise Ratio (SNR)', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'snr_analysis.png'), 0.3, 1.1, width=12.5)
tf = add_textbox(slide, 0.5, 6.2, 12, 1,
                 'Top: raw signal | Middle: kept (signal) | Bottom: removed (noise) — FIR effectively separates speech from noise',
                 font_size=16, color=GRAY)

# ============================================================
# SLIDE 11: Pre-emphasis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, "Pre-emphasis: y'[n] = y[n] - 0.97 * y[n-1]", font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'preemphasis_effect.png'), 0.3, 1.2, width=12.5)
tf = add_textbox(slide, 0.5, 5.5, 12, 1.5,
                 'Left: After FIR only — high frequencies still weak\n'
                 'Right: After FIR + Pre-emphasis — spectrum flattened, high frequencies boosted for better MFCC extraction',
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 12: Speaker spectrum comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Spectrum Comparison Across Speakers', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'speakers_spectrum.png'), 0.8, 1.0, width=11.5)
tf = add_textbox(slide, 0.5, 6.4, 12, 0.8,
                 'Each speaker has a unique spectral fingerprint — different formant positions → basis for MFCC features',
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 13: Feature Engineering — Pipeline A vs B
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Feature Engineering', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.3, 1.5, 6, 5, 'Pipeline A — Baseline (6-dim)', font_size=24, bold=True, color=FPT)
add_para(tf, '', font_size=8)
add_para(tf, '1. RMS Energy (mean) — loudness', font_size=18)
add_para(tf, '2. RMS Energy (std) — energy variation', font_size=18)
add_para(tf, '3. ZCR (mean) — zero crossing rate', font_size=18)
add_para(tf, '4. ZCR (std) — ZCR variation', font_size=18)
add_para(tf, '5. Mean |amplitude|', font_size=18)
add_para(tf, '6. Std amplitude', font_size=18)
add_para(tf, '', font_size=8)
add_para(tf, 'Only time-domain — NO frequency info!', font_size=18, bold=True, color=RGBColor(0xE7, 0x4C, 0x3C))

tf2 = add_textbox(slide, 6.8, 1.5, 6, 5, 'Pipeline B — MFCC (26-dim)', font_size=24, bold=True, color=FPT)
add_para(tf2, '', font_size=8)
add_para(tf2, '1. Frame signal (512 samples, hop=256)', font_size=18)
add_para(tf2, '2. Hamming window each frame', font_size=18)
add_para(tf2, '3. FFT → power spectrum', font_size=18)
add_para(tf2, '4. Mel filterbank (mimics human ear)', font_size=18)
add_para(tf2, '5. Log of filterbank energies', font_size=18)
add_para(tf2, '6. DCT → 13 MFCCs per frame', font_size=18)
add_para(tf2, '', font_size=8)
add_para(tf2, 'Final: mean + std of 13 MFCCs = 26-dim', font_size=18, bold=True, color=RGBColor(0x27, 0xAE, 0x60))

# ============================================================
# SLIDE 14: MFCC Heatmap + Feature Scatter
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Feature Visualization', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'mfcc_heatmap.png'), 0.2, 1.1, width=6.3)
add_image(slide, os.path.join(FIGURES, 'feature_scatter.png'), 6.7, 1.1, width=6.3)
tf = add_textbox(slide, 0.5, 6.0, 12, 1,
                 'Left: MFCC heatmap (13 coefficients over time) | Right: Feature scatter — speakers form distinct clusters',
                 font_size=18, color=GRAY)

# ============================================================
# SLIDE 15: SVM Classifier
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'SVM with RBF Kernel', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.5, 6, 5, 'Why SVM?', font_size=26, bold=True)
add_para(tf, '', font_size=8)
add_para(tf, '- Works well with small datasets (125 samples)', font_size=20)
add_para(tf, '- Effective in high-dimensional space (26-dim)', font_size=20)
add_para(tf, '- Clear decision boundaries', font_size=20)
add_para(tf, '', font_size=12)
add_para(tf, 'K(xi, xj) = exp(-gamma * ||xi - xj||^2)', font_size=20, bold=True, color=FPT)

tf2 = add_textbox(slide, 7, 1.5, 5.5, 5, 'Training Setup', font_size=26, bold=True)
add_para(tf2, '', font_size=8)
add_para(tf2, 'Pipeline: StandardScaler → SVM', font_size=20)
add_para(tf2, '', font_size=8)
add_para(tf2, 'Hyperparameter Search:', font_size=20, bold=True)
add_para(tf2, '  C:  0.1, 1, 10, 100', font_size=20)
add_para(tf2, '  gamma:  scale, auto, 0.001, 0.01', font_size=20)
add_para(tf2, '', font_size=8)
add_para(tf2, 'GridSearchCV (3-fold inner)', font_size=20)
add_para(tf2, '5-fold Stratified CV (seed=42)', font_size=20)

# ============================================================
# SLIDE 16: CV Scores
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Cross-Validation Results', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'cv_scores.png'), 0.8, 1.1, width=11.5)
tf = add_textbox(slide, 0.5, 6.0, 12, 1,
                 'Pipeline A: 56.3% ± 2.8%  |  Pipeline B: 97.0% ± 3.6%  |  Gap: +40.7 percentage points',
                 font_size=20, bold=True, color=FPT, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 17: Results table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Experimental Results', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.5, 12, 4, '', font_size=20)
rows = [
    ('Pipeline', 'Best C', 'Best gamma', 'Mean Accuracy', '95% CI', 'F1-score'),
    ('A — Baseline', '1', 'scale', '56.3% ± 2.8%', '[52.4%, 60.1%]', '0.690'),
    ('B — DSP-enhanced', '1', 'scale', '97.0% ± 3.6%', '[92.0%, 102.1%]', '1.000'),
]
for i, row in enumerate(rows):
    line = '    '.join(f'{col:<20}' for col in row)
    if i == 0:
        tf.paragraphs[0].text = line
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
    else:
        p = add_para(tf, line, font_size=18, space_before=Pt(12))
        if i == 2:
            p.font.color.rgb = FPT
            p.font.bold = True

add_para(tf, '', font_size=14)
add_para(tf, 'Paired t-test:  t = -34.79,  p = 0.000004', font_size=24, bold=True, space_before=Pt(24))
add_para(tf, 'p << 0.05  →  HIGHLY STATISTICALLY SIGNIFICANT', font_size=22, bold=True, color=FPT)

# ============================================================
# SLIDE 18: Confusion Matrices
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Confusion Matrices', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'confusion_matrices.png'), 1.5, 1.1, width=10)
tf = add_textbox(slide, 0.5, 5.8, 12, 1,
                 'Pipeline A: many misclassifications | Pipeline B: near-perfect diagonal — correct predictions',
                 font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 19: ROC Curves
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'ROC Curves (One-vs-Rest)', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'roc_curves.png'), 1.5, 1.1, width=10)
tf = add_textbox(slide, 0.5, 5.8, 12, 1,
                 'Pipeline B: AUC close to 1.0 for all speakers — excellent discrimination ability',
                 font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 20: Comparison bar chart
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide, height=0.9)
add_textbox(slide, 0.5, 0.1, 12, 0.7, 'Performance Comparison', font_size=32, bold=True, color=WHITE)
add_image(slide, os.path.join(FIGURES, 'comparison_bar.png'), 2.5, 1.2, width=8)
tf = add_textbox(slide, 0.5, 5.8, 12, 1,
                 'Pipeline B outperforms A on all metrics: Accuracy, Precision, Recall, F1-score',
                 font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 21: Discussion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Why Pipeline B Wins', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.5, 12, 5, '', font_size=20)
tf.paragraphs[0].text = '1. Feature Richness'
tf.paragraphs[0].font.size = Pt(26)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = FPT
add_para(tf, '   26-dim MFCCs capture vocal tract shape', font_size=20)
add_para(tf, '   vs. 6-dim time features — only amplitude info', font_size=20, color=GRAY)
add_para(tf, '', font_size=12)

add_para(tf, '2. Noise Suppression', font_size=26, bold=True, color=FPT, space_before=Pt(16))
add_para(tf, '   FIR bandpass + pre-emphasis isolate speech frequencies', font_size=20)
add_para(tf, '   → improved SNR before feature extraction', font_size=20, color=GRAY)
add_para(tf, '', font_size=12)

add_para(tf, '3. Massive Accuracy Gap', font_size=26, bold=True, color=FPT, space_before=Pt(16))
add_para(tf, '   +40.7 percentage points (56.3% → 97.0%)', font_size=22, bold=True)
add_para(tf, '   Computational cost: <1 ms per 3-second clip', font_size=20, color=GRAY)

# ============================================================
# SLIDE 22: Limitations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Limitations', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.5, 12, 5, '', font_size=20)
items = [
    ('Small dataset', '5 speakers, 125 samples — may not generalize to larger populations'),
    ('Clean recordings only', 'Trained in quiet; noisy environments (classroom) degrade mic predictions'),
    ('Closed-set', 'Cannot reject unknown speakers — assumes test speaker is known'),
    ('Single classifier', 'Only SVM tested — RF, neural networks may respond differently'),
    ('Fixed filter design', '300–3400 Hz passband not optimized per dataset'),
]
tf.paragraphs[0].text = ''
for i, (title, desc) in enumerate(items):
    p = add_para(tf, f'  {title}', font_size=22, bold=True, space_before=Pt(16))
    add_para(tf, f'      {desc}', font_size=18, color=GRAY)

# ============================================================
# SLIDE 23: Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_orange_bar(slide)
add_textbox(slide, 0.5, 0.2, 12, 0.8, 'Conclusion', font_size=36, bold=True, color=WHITE)

tf = add_textbox(slide, 0.5, 1.5, 12, 3.5, 'Key Findings:', font_size=28, bold=True)
add_para(tf, '', font_size=8)
add_para(tf, '1.  DSP preprocessing is essential — 40.7 pp accuracy gain', font_size=22)
add_para(tf, '     Confirmed by paired t-test (p = 0.000004)', font_size=18, color=GRAY)
add_para(tf, '', font_size=8)
add_para(tf, '2.  MFCC features capture speaker-specific vocal tract characteristics', font_size=22)
add_para(tf, '     Far better than basic time-domain descriptors (RMS, ZCR)', font_size=18, color=GRAY)
add_para(tf, '', font_size=8)
add_para(tf, '3.  Handcrafted DSP + SVM achieves 97% on small dataset', font_size=22)
add_para(tf, '     No deep learning required!', font_size=18, color=FPT, bold=True)

tf2 = add_textbox(slide, 0.5, 5.5, 12, 1.5, 'Future Work:', font_size=24, bold=True)
add_para(tf2, '  - More speakers, noisy conditions, noise reduction for real-time mic', font_size=18, color=GRAY)
add_para(tf2, '  - Integrate DSP front-end with neural network architectures', font_size=18, color=GRAY)

# ============================================================
# SLIDE 24: Live Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FPT)
add_textbox(slide, 0.5, 2.0, 12, 1.5, 'Live Demo',
            font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 3.8, 12, 1, 'Streamlit App — Speaker Identification',
            font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 4.8, 12, 1, 'Upload WAV file or Record from mic',
            font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 25: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FPT)
add_textbox(slide, 0.5, 2.2, 12, 1.5, 'Thank You!',
            font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 4.0, 12, 1, 'Questions?',
            font_size=32, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f'Saved: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
